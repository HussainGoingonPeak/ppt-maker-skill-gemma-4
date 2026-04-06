from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================================
# PPT MAKER - Professional PowerPoint Generator
# ============================================================

class PPTMaker:
    def __init__(self, theme="modern"):
        self.prs = Presentation()
        self.theme = self._get_theme(theme)
        self._setup_slide_size()
    
    def _get_theme(self, theme_name):
        """Define color schemes and styling"""
        themes = {
            "modern": {
                "primary": (30, 58, 138),      # Deep Blue
                "secondary": (59, 130, 246),   # Bright Blue
                "accent": (245, 158, 11),      # Amber
                "bg": (255, 255, 255),         # White
                "text": (31, 41, 55),          # Dark Gray
                "light_bg": (243, 244, 246)    # Light Gray
            },
            "dark": {
                "primary": (17, 24, 39),       # Near Black
                "secondary": (55, 65, 81),     # Dark Gray
                "accent": (16, 185, 129),      # Emerald
                "bg": (249, 250, 251),         # Off White
                "text": (17, 24, 39),          # Near Black
                "light_bg": (229, 231, 235)    # Gray
            },
            "creative": {
                "primary": (124, 58, 237),     # Violet
                "secondary": (236, 72, 153),   # Pink
                "accent": (251, 191, 36),      # Amber
                "bg": (255, 255, 255),         # White
                "text": (31, 41, 55),          # Dark Gray
                "light_bg": (254, 243, 199)    # Light Yellow
            }
        }
        return themes.get(theme_name, themes["modern"])
    
    def _setup_slide_size(self):
        """Set 16:9 widescreen format"""
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def _add_shape_background(self, slide, shape_type, left, top, width, height, color):
        """Add colored shape as background element"""
        shape = slide.shapes.add_shape(
            shape_type, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(*color)
        shape.line.fill.background()
        return shape
    
    def _style_text_frame(self, text_frame, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        """Apply consistent text styling"""
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = align
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.name = "Calibri"
                if color:
                    run.font.color.rgb = RgbColor(*color)
                else:
                    run.font.color.rgb = RgbColor(*self.theme["text"])
    
    def create_title_slide(self, title, subtitle=None, presenter=None):
        """Create opening title slide with modern design"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add decorative header bar
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2),
            self.theme["primary"]
        )
        
        # Add accent line
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(1.2),
            Inches(13.333), Inches(0.1),
            self.theme["accent"]
        )
        
        # Title textbox
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = title
        self._style_text_frame(tf, font_size=44, bold=True, color=self.theme["primary"], align=PP_ALIGN.CENTER)
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
            )
            tf = sub_box.text_frame
            tf.text = subtitle
            self._style_text_frame(tf, font_size=24, color=self.theme["secondary"], align=PP_ALIGN.CENTER)
        
        # Presenter info
        if presenter:
            pres_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6), Inches(12.333), Inches(0.8)
            )
            tf = pres_box.text_frame
            tf.text = f"Presented by: {presenter}"
            self._style_text_frame(tf, font_size=16, color=self.theme["text"], align=PP_ALIGN.CENTER)
        
        return slide
    
    def add_content_slide(self, title, content_items, subtitle=None):
        """Add standard content slide with bullet points"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Header background
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.3),
            self.theme["primary"]
        )
        
        # Title on slide
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        tf = title_box.text_frame
        tf.text = title
        self._style_text_frame(tf, font_size=32, bold=True, color=(255, 255, 255), align=PP_ALIGN.LEFT)
        
        # Content area background
        self._add_shape_background(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.6),
            Inches(12.333), Inches(5.5),
            self.theme["light_bg"]
        )
        
        # Content textbox
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.9), Inches(11.733), Inches(5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {item}"
            p.level = 0
            p.font.size = Pt(20)
            p.font.name = "Calibri"
            p.font.color.rgb = RgbColor(*self.theme["text"])
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(self, title, left_title, left_items, right_title, right_items):
        """Create side-by-side content layout"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Header
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2),
            self.theme["primary"]
        )
        
        # Main title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
        )
        tf = title_box.text_frame
        tf.text = title
        self._style_text_frame(tf, font_size=28, bold=True, color=(255, 255, 255))
        
        # Left column
        left_bg = self._add_shape_background(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.5),
            Inches(6), Inches(5.5),
            (255, 255, 255)
        )
        
        left_title_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7), Inches(5.6), Inches(0.6)
        )
        tf = left_title_box.text_frame
        tf.text = left_title
        self._style_text_frame(tf, font_size=22, bold=True, color=self.theme["primary"])
        
        left_content = slide.shapes.add_textbox(
            Inches(0.7), Inches(2.4), Inches(5.6), Inches(4.4)
        )
        tf = left_content.text_frame
        tf.word_wrap = True
        for item in left_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        
        # Right column
        right_bg = self._add_shape_background(
            slide, MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.8), Inches(1.5),
            Inches(6), Inches(5.5),
            self.theme["light_bg"]
        )
        
        right_title_box = slide.shapes.add_textbox(
            Inches(7), Inches(1.7), Inches(5.6), Inches(0.6)
        )
        tf = right_title_box.text_frame
        tf.text = right_title
        self._style_text_frame(tf, font_size=22, bold=True, color=self.theme["secondary"])
        
        right_content = slide.shapes.add_textbox(
            Inches(7), Inches(2.4), Inches(5.6), Inches(4.4)
        )
        tf = right_content.text_frame
        tf.word_wrap = True
        for item in right_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        
        return slide
    
    def add_section_divider(self, section_title, section_number=None):
        """Create section transition slide"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Full background
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(7.5),
            self.theme["primary"]
        )
        
        # Accent decoration
        self._add_shape_background(
            slide, MSO_SHAPE.OVAL,
            Inches(10), Inches(-2),
            Inches(5), Inches(5),
            self.theme["secondary"]
        )
        
        # Section number
        if section_number:
            num_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.5), Inches(2), Inches(1.5)
            )
            tf = num_box.text_frame
            tf.text = f"0{section_number}" if section_number < 10 else str(section_number)
            self._style_text_frame(tf, font_size=72, bold=True, color=self.theme["accent"])
        
        # Section title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.text = section_title
        self._style_text_frame(tf, font_size=48, bold=True, color=(255, 255, 255))
        
        # Accent line
        self._add_shape_background(
            slide, MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(5.5),
            Inches(3), Inches(0.15),
            self.theme["accent"]
        )
        
        return slide
    
    def add_image_slide(self, title, image_path, caption=None, full_bleed=False):
        """Add slide with image"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if not full_bleed:
            # Header
            self._add_shape_background(
                slide, MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(1.2),
                self.theme["primary"]
            )
            
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.7)
            )
            tf = title_box.text_frame
            tf.text = title
            self._style_text_frame(tf, font_size=28, bold=True, color=(255, 255, 255))
            
            # Image area
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    Inches(1), Inches(1.8),
                    width=Inches(11.333)
                )
            
            if caption:
                cap_box = slide.shapes.add_textbox(
                    Inches(1), Inches(6.5), Inches(11.333), Inches(0.6)
                )
                tf = cap_box.text_frame
                tf.text = caption
                self._style_text_frame(tf, font_size=14, color=self.theme["text"], align=PP_ALIGN.CENTER)
        else:
            # Full background image style
            if os.path.exists(image_path):
                slide.shapes.add_picture(
                    image_path,
                    Inches(0), Inches(0),
                    height=Inches(7.5)
                )
        
        return slide
    
    def save(self, filename="presentation.pptx"):
        """Save presentation to file"""
        output_path = f"/mnt/kimi/output/{filename}"
        self.prs.save(output_path)
        return output_path


# ============================================================
# QUICK START FUNCTIONS
# ============================================================

def quick_ppt(title, slides_data, theme="modern", filename="output.pptx"):
    """
    Quick creation function
    slides_data: list of dicts with 'type', 'title', 'content', etc.
    """
    maker = PPTMaker(theme=theme)
    
    # Title slide
    maker.create_title_slide(
        title=title,
        subtitle=slides_data[0].get('subtitle') if slides_data else None
    )
    
    # Content slides
    for slide in slides_data[1:]:
        slide_type = slide.get('type', 'content')
        
        if slide_type == 'content':
            maker.add_content_slide(
                title=slide['title'],
                content_items=slide['content']
            )
        elif slide_type == 'two_column':
            maker.add_two_column_slide(
                title=slide['title'],
                left_title=slide['left_title'],
                left_items=slide['left_content'],
                right_title=slide['right_title'],
                right_items=slide['right_content']
            )
        elif slide_type == 'section':
            maker.add_section_divider(
                section_title=slide['title'],
                section_number=slide.get('number')
            )
        elif slide_type == 'image':
            maker.add_image_slide(
                title=slide['title'],
                image_path=slide['image_path'],
                caption=slide.get('caption')
            )
    
    return maker.save(filename)


# ============================================================
# EXAMPLE USAGE
# ============================================================

if __name__ == "__main__":
    # Example: Create a sample presentation
    slides = [
        {
            'type': 'title',
            'subtitle': 'Q4 2024 Business Review'
        },
        {
            'type': 'section',
            'title': 'Executive Summary',
            'number': 1
        },
        {
            'type': 'content',
            'title': 'Key Highlights',
            'content': [
                'Revenue growth of 25% year-over-year',
                'Expanded into 3 new markets',
                'Launched 2 major product features',
                'Customer satisfaction score: 94%'
            ]
        },
        {
            'type': 'two_column',
            'title': 'Market Analysis',
            'left_title': 'Opportunities',
            'left_content': [
                'Emerging AI market segment',
                'Strategic partnerships available',
                'Untapped enterprise clients'
            ],
            'right_title': 'Challenges',
            'right_content': [
                'Increased competition',
                'Supply chain constraints',
                'Talent acquisition costs'
            ]
        },
        {
            'type': 'section',
            'title': 'Next Steps',
            'number': 2
        }
    ]
    
    output_file = quick_ppt(
        title="Quarterly Business Review",
        slides_data=slides,
        theme="modern",
        filename="business_review.pptx"
    )
    
    print(f"Presentation saved to: {output_file}")
      
